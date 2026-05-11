"""
Generate ADHD Productivity MVP PowerPoint Presentation
Professional 10-slide presentation with clean design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_COLOR = RGBColor(41, 98, 255)      # Professional Blue
ACCENT_COLOR = RGBColor(108, 92, 231)     # Purple accent
DARK_TEXT = RGBColor(44, 62, 80)          # Dark gray
LIGHT_BG = RGBColor(240, 242, 247)        # Light blue background
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "May 5, 2026 | Rishi Kumar"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, accent_color=PRIMARY_COLOR):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = accent_color
    header.line.color.rgb = accent_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
        p.font.bold = False
    
    return slide

def add_two_column_slide(prs, title, left_bullets, right_bullets, accent_color=PRIMARY_COLOR):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header bar
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = accent_color
    header.line.color.rgb = accent_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
add_title_slide(
    prs,
    "ADHD Productivity MVP",
    "AI-Based ADHD Productivity Support System"
)

# ============================================================================
# SLIDE 2: PROJECT OVERVIEW
# ============================================================================
add_content_slide(
    prs,
    "Project Overview 🎯",
    [
        "🧠 Mission: Intelligent ADHD productivity support system",
        "",
        "Problem Statement:",
        "  • ADHD affects 5-10% of global population",
        "  • Productivity challenges: distractions, time management",
        "  • Mental health struggles: stress, anxiety, depression",
        "  • Need for personalized, real-time support",
        "",
        "Solution: AI-powered analytics + behavioral insights"
    ],
    PRIMARY_COLOR
)

# ============================================================================
# SLIDE 3: OBJECTIVES
# ============================================================================
add_content_slide(
    prs,
    "Key Objectives 🎯",
    [
        "✅ Detect ADHD risk factors using behavioral analysis",
        "",
        "✅ Monitor real-time mental health and stress levels",
        "",
        "✅ Generate personalized productivity interventions",
        "",
        "✅ Provide interactive AI chatbot support",
        "",
        "✅ Ensure efficient, scalable ML inference",
        "",
        "✅ Achieve production-ready reliability and performance"
    ],
    ACCENT_COLOR
)

# ============================================================================
# SLIDE 4: SYSTEM ARCHITECTURE
# ============================================================================
add_content_slide(
    prs,
    "System Architecture 🏗️",
    [
        "🎨 Frontend Layer: Streamlit UI",
        "  • Chat interface, real-time score display",
        "  • Session management & persistence",
        "",
        "⚙️ API Layer: FastAPI (15+ endpoints)",
        "  • /calculate_scores, /chat, /get_interventions",
        "",
        "🚀 Optimized Inference Layer",
        "  • Model caching (1343x improvement)",
        "  • Batch processing & prediction caching",
        "",
        "🤖 ML Models Layer: 4 Production Models",
        "  • CatBoost, RandomForest, TF-IDF + Logistic Regression"
    ],
    PRIMARY_COLOR
)

# ============================================================================
# SLIDE 5: CORE FEATURES
# ============================================================================
add_two_column_slide(
    prs,
    "Core Features ⚡",
    [
        "🧬 ADHD Risk Assessment",
        "  • Behavioral pattern analysis",
        "  • Real-time risk scoring (0-100)",
        "",
        "📊 Productivity Scoring",
        "  • Task completion metrics",
        "  • Focus & output quality"
    ],
    [
        "🧘 Mental Health Monitoring",
        "  • Multi-level stress detection",
        "  • Depression risk assessment",
        "",
        "💬 Interactive Chatbot",
        "  • ADHD-specific guidance",
        "  • Offline capability"
    ],
    ACCENT_COLOR
)

# ============================================================================
# SLIDE 6: MACHINE LEARNING MODELS
# ============================================================================
add_content_slide(
    prs,
    "ML Models Deployed 🤖",
    [
        "1️⃣ ADHD Risk Model (CatBoost) - 345 KB",
        "   → Predicts ADHD risk probability",
        "",
        "2️⃣ Productivity Model (CatBoost) - 178 KB",
        "   → Estimates productivity score",
        "",
        "3️⃣ Mental Health NLP (TF-IDF + LogReg) - 454 KB",
        "   → Text-based stress analysis",
        "",
        "4️⃣ Student Depression (RandomForest) - 696 KB",
        "   → Student mental health prediction"
    ],
    PRIMARY_COLOR
)

# ============================================================================
# SLIDE 7: PERFORMANCE & OPTIMIZATION
# ============================================================================
add_two_column_slide(
    prs,
    "Performance & Optimization ⚡",
    [
        "🚀 Optimization Techniques:",
        "  • Model caching (1343x improvement)",
        "  • Batch prediction (1.7-2x)",
        "  • Prediction caching (instant)",
        "  • Feature alignment (30-40%)",
        "  • Thread safety (stable)",
        "  • Memory efficiency (bounded)"
    ],
    [
        "📈 Performance Results:",
        "  • Cold startup: 5x faster",
        "  • API response: 50-100ms",
        "  • Throughput: 80-120 req/s",
        "  • Memory: 200-300MB bounded",
        "",
        "✅ Total Improvement: 10-500x"
    ],
    ACCENT_COLOR
)

# ============================================================================
# SLIDE 8: KEY ACHIEVEMENTS
# ============================================================================
add_content_slide(
    prs,
    "Key Achievements 🏆",
    [
        "✅ Quality: 100% test pass rate (32/32 tests passing)",
        "",
        "✅ Deployment: 4 ML models production-ready",
        "",
        "✅ Scale: 15+ API endpoints functional",
        "",
        "✅ Optimization: 6 optimization techniques integrated",
        "",
        "✅ Documentation: 15+ comprehensive guides (82+ KB)",
        "",
        "✅ Performance: 10-500x improvement across all models"
    ],
    PRIMARY_COLOR
)

# ============================================================================
# SLIDE 9: APPLICATIONS & FUTURE SCOPE
# ============================================================================
add_two_column_slide(
    prs,
    "Applications & Future Scope 🚀",
    [
        "Current Applications:",
        "  • Student mental health monitoring",
        "  • Workplace productivity tracking",
        "  • Clinical ADHD assessment support",
        "  • Remote mental health services"
    ],
    [
        "Future Enhancements:",
        "  • Wearable data integration (heart rate, sleep)",
        "  • Advanced NLP (BERT, transformers)",
        "  • Mobile app (React Native)",
        "  • Telehealth professional integration",
        "  • Multi-language support"
    ],
    ACCENT_COLOR
)

# ============================================================================
# SLIDE 10: CONCLUSION
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PRIMARY_COLOR

# Main message
main_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
main_frame = main_box.text_frame
main_frame.word_wrap = True
p = main_frame.paragraphs[0]
p.text = "Project Impact & Status"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Impact points
impact_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(2.5))
impact_frame = impact_box.text_frame
impact_frame.word_wrap = True

points = [
    "🎯 Complete system ready for production deployment",
    "📊 Proven performance: 10-500x optimization",
    "✅ All systems operational & thoroughly tested",
    "",
    "🟢 READY FOR IMMEDIATE DEPLOYMENT & USER ROLLOUT"
]

for i, point in enumerate(points):
    if i == 0:
        p = impact_frame.paragraphs[0]
    else:
        p = impact_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(22) if i == 4 else Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True if i == 4 else False
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(4)
    p.space_after = Pt(4)

# Save presentation
output_path = r"d:\ADHD_Productivity_MVP\ADHD_Productivity_MVP_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
